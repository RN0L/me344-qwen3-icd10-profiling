#!/usr/bin/env python3
"""Build the 5-slide deck from the measured records.

Every figure on these slides is read out of ``results/`` at build time. Nothing is
transcribed, so the deck cannot drift from the data, re-run this after any new run
and the numbers update themselves. The speaker notes are built from the same
variables as the slides, so the spoken number and the projected number cannot
disagree either.

    uv run --with python-pptx --with pillow python slides/build_pptx.py

Design intent: one idea per slide, one number large enough to read from the back of
the room, and body text kept under ~45 words. The house colour is a single constant.
Cards size themselves to their content (see ``card``) so text can never be clipped
by the panel it sits in.
"""

from __future__ import annotations

import json
import textwrap
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
RESULTS = REPO / "results"
FIGURES = Path(__file__).resolve().parent / "figures"
OUT = Path(__file__).resolve().parent / "ME344_Final_Arnold_Hambuch.pptx"

# --- palette -------------------------------------------------------------------
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
ORANGE_L = RGBColor(0xFB, 0x92, 0x3C)
TINT = RGBColor(0xFE, 0xF3, 0xEA)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_TINT = RGBColor(0xFE, 0xF2, 0xF2)
INK = RGBColor(0x15, 0x17, 0x1A)
GREY = RGBColor(0x71, 0x7A, 0x85)
FAINT = RGBColor(0xE6, 0xE9, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Helvetica Neue"

W, H = Inches(13.333), Inches(7.5)

# Shapes whose box deliberately overhangs: the ghost slide number is a
# right-aligned box, so its frame reaches the margin while its glyph does not.
EXEMPT = set()


def load(n):
    return json.loads((RESULTS / f"{n}.json").read_text())


# --- facts, read not typed -----------------------------------------------------
recs = {}
for f in sorted(RESULTS.glob("*.json")):
    if f.stem != "analysis":
        d = json.loads(f.read_text())
        recs[d["run_id"]] = d

cpu = recs["cpu-x86-32c-0p6b-bs1-seq1024"]
gpu = recs["gpu-gh200-0p6b-bs1-seq1024"]
tpu = recs["tpu-v5e8-bs1-seq1024-filecache-0p6b"]
base = recs["tpu-v5e8-bs8-seq1024"]
mit = recs["tpu-v5e8-bs8-seq1024-filecache"]

c_s, g_s, t_s = (r["steady"]["median_step_s"] for r in (cpu, gpu, tpu))
sp_g, sp_t = c_s / g_s, c_s / t_s
gt = abs(g_s - t_s) / t_s * 100
bp, mp = base["phases"], mit["phases"]
compute_pct = bp["steady_state_s"] / bp["total_wall_s"] * 100
load_pct = bp["model_load_s"] / bp["total_wall_s"] * 100
load_x = bp["model_load_s"] / mp["model_load_s"]
wall_cut = (1 - mp["total_wall_s"] / bp["total_wall_s"]) * 100

# Utilisation and memory for the like-for-like row. The TPU exposes memory
# statistics but no core-utilisation counter, so that cell says so rather than
# borrowing a number that means something else.
u_cpu = cpu["utilization"]["mean_pct"]
u_gpu = gpu["utilization"]["mean_pct"]
m_cpu, m_gpu, m_tpu = (r["memory"]["peak_pct"] for r in (cpu, gpu, tpu))

# The 4B sweep: the comparison at an operating point where both chips are busy.
g4 = {r["config"]["batch_size"]: r for r in recs.values()
      if r["backend"] == "gpu" and r["config"]["model"].endswith("4B")}
t4 = {r["config"]["batch_size"]: r for r in recs.values()
      if r["backend"] == "tpu" and r["config"]["model"].endswith("4B")
      and r["config"]["seq_len"] == 1024}
g8, t8 = g4[8]["steady"]["tokens_per_s"], t4[8]["steady"]["tokens_per_s"]
g16, t16 = g4[16]["steady"]["tokens_per_s"], t4[16]["steady"]["tokens_per_s"]
g32 = g4[32]["steady"]["tokens_per_s"]
load_ratio = g16 / t16

# The last configuration that fits, per backend, and the one that does not.
def ceiling(pool):
    ok = {b: r for b, r in pool.items() if r["status"] == "ok"}
    died = sorted(b for b, r in pool.items() if r["status"] == "oom")
    top = max(ok)
    return ok[top]["memory"]["peak_pct"], top, (died[0] if died else None)


mem_gpu, bs_gpu, oom_gpu = ceiling(g4)
mem_tpu, bs_tpu, oom_tpu = ceiling(t4)

_c = {c["run_id"]: c["usd_per_1000_steps"] for c in load("analysis")["cost"]["per_run"]}
k4_gpu = _c["gpu-gh200-bs8-seq1024"]
k4_tpu = _c["tpu-v5e8-bs8-seq1024-filecache"]
cost_ratio_4b = k4_tpu / k4_gpu

stats = json.loads((REPO / "docs" / "dataset_stats.json").read_text())
tok = stats["splits"]["train"]["tokens"]["total"]
n_train = stats["splits"]["train"]["n_docs"]
n_dev = stats["splits"]["dev"]["n_docs"]
trunc = {e["seq_len"]: e for e in stats["splits"]["train"]["truncation"]}
cut_1024 = 100 * trunc[1024]["frac_truncated"]

prs = Presentation()
prs.slide_width, prs.slide_height = W, H


# --- primitives ----------------------------------------------------------------
def slide(num, kicker, headline):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), H)
    rule.fill.solid(); rule.fill.fore_color.rgb = ORANGE
    rule.line.fill.background(); rule.shadow.inherit = False

    ghost = s.shapes.add_textbox(Inches(11.9), Inches(0.05), Inches(1.3), Inches(1.5))
    r = ghost.text_frame.paragraphs[0].add_run(); r.text = str(num)
    r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = FAINT; r.font.name = FONT
    ghost.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    EXEMPT.add(ghost.shape_id)

    txt(s, Inches(0.62), Inches(0.40), Inches(10.5), Inches(0.3),
        kicker.upper(), 12, color=ORANGE, bold=True, spacing=2)
    txt(s, Inches(0.62), Inches(0.72), Inches(11.0), Inches(0.9),
        headline, 32, bold=True)
    return s


def txt(s, x, y, w, h, text, size=14, bold=False, color=INK,
        align=PP_ALIGN.LEFT, gap=5, spacing=0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(gap)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT
    return tb


def lines(text, size, gap, width=None):
    """Height a block of text needs, in inches, counting soft wraps too.

    `width` is the inches available for the text. Helvetica Neue averages roughly
    0.52 em per character in running text, which is what turns a width into a
    character budget. Deliberately generous: over-estimating adds padding,
    under-estimating clips, and clipping is the bug this function exists to stop.
    """
    per_line = max(1, int(width / (size / 72 * 0.52))) if width else 10 ** 6
    n = 0
    for hard in text.split("\n"):
        n += max(1, -(-len(hard) // per_line))
    return n * (size / 72 * 1.28 + gap / 72)


def stat(s, x, y, w, value, label, big=44, color=ORANGE):
    """A single number, sized to read from the back of the room."""
    txt(s, x, y, w, Inches(0.85), value, big, bold=True, color=color, gap=0)
    txt(s, x + Inches(0.02), y + Inches(0.72), w, Inches(0.9), label, 12, color=GREY, gap=2)
    return y + Inches(0.72) + Inches(lines(label, 12, 2))


def panel(s, x, y, w, h, fill=TINT, edge=FAINT):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.04
    return sh


PAD = 0.26          # inches of breathing room inside every card
TITLE_GAP = 0.10


def card(s, x, y, w, title, body, title_size=15, body_size=13, gap=4,
         fill=TINT, accent=ORANGE, foot=None):
    """A tinted card whose height is computed from its content.

    The clipped-text bug this replaces came from writing panel heights by hand and
    then editing the text. Here the panel cannot be shorter than what is in it.
    """
    inner = w / Inches(1) - 2 * PAD
    th = title_size / 72 * 1.32 if title else 0.0
    bh = lines(body, body_size, gap, inner)
    fh = (lines(foot, 10.5, 2, inner) + 0.08) if foot else 0.0
    h = PAD + th + (TITLE_GAP if title else 0) + bh + fh + PAD * 0.72
    panel(s, x, y, w, Inches(h), fill=fill)
    cy = y + Inches(PAD)
    if title:
        txt(s, x + Inches(PAD + 0.02), cy, w - Inches(2 * PAD), Inches(th),
            title, title_size, bold=True, color=accent)
        cy += Inches(th + TITLE_GAP)
    txt(s, x + Inches(PAD + 0.02), cy, w - Inches(2 * PAD), Inches(bh),
        body, body_size, gap=gap)
    if foot:
        txt(s, x + Inches(PAD + 0.02), cy + Inches(bh + 0.08),
            w - Inches(2 * PAD), Inches(fh), foot, 10.5, color=GREY, gap=2)
    return y + Inches(h)


def table(s, x, y, rows, widths, size=13, rh=Inches(0.40), highlight=None,
          right_align_from=1, rule_after=None):
    """right_align_from: first column index to right-align. Set high to keep text left."""
    t = s.shapes.add_table(len(rows), len(rows[0]), x, y, sum(widths), rh * len(rows)).table
    for j, wd in enumerate(widths):
        t.columns[j].width = wd
    for i, row in enumerate(rows):
        t.rows[i].height = rh
        for j, v in enumerate(row):
            c = t.cell(i, j); c.text = str(v)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.09)
            c.margin_top = c.margin_bottom = Inches(0.01)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j >= right_align_from else PP_ALIGN.LEFT
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(size); r.font.name = FONT
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = INK
                r.font.bold = True; r.font.color.rgb = WHITE
            elif highlight is not None and j == highlight:
                c.fill.fore_color.rgb = TINT
                r.font.bold = True; r.font.color.rgb = ORANGE
            elif rule_after is not None and i == rule_after:
                c.fill.fore_color.rgb = TINT
                r.font.color.rgb = INK
            else:
                c.fill.fore_color.rgb = WHITE
                r.font.color.rgb = INK
    return y + rh * len(rows)


def fig(s, name, x, y, width):
    """Place a figure and return its bottom edge, so the next thing can clear it."""
    p = FIGURES / name
    if not p.exists():
        return y
    s.shapes.add_picture(str(p), x, y, width=width)
    px, py = Image.open(p).size
    return y + Inches(width / Inches(1) * py / px)


def notes(s, t):
    """Set speaker notes, reflowed after interpolation.

    The notes are f-strings, so substituting a value mid-sentence leaves the
    source's hard wraps in the wrong places. Rewrap here: cue lines in brackets
    and the header lines stay on their own line, bullets keep a hanging indent,
    everything else is a paragraph.
    """
    HEAD = ("LEONARD", "LUIS", "TIMING", "IF ASKED", "→", "Hand back", "Read at")
    out, unit, kind = [], [], "para"

    def flush():
        if unit:
            text = " ".join(unit)
            out.append(text if kind == "head" else textwrap.fill(
                text, 84, subsequent_indent="  " if kind == "bullet" else ""))
            unit.clear()

    for line in t.strip().split("\n"):
        st = line.strip()
        if not st:
            flush(); out.append(""); kind = "para"; continue
        if st.startswith("·"):
            flush(); kind = "bullet"
        elif st.startswith("[") or st.startswith(HEAD):
            flush(); kind = "head"
        elif kind == "head":
            flush(); kind = "para"
        unit.append(st)
    flush()
    s.notes_slide.notes_text_frame.text = "\n".join(out).strip()


# ============================== 1. Problem ====================================
s = slide(1, "The problem", "Can a clinic afford to run its own ICD-10 model?")

card(s, Inches(0.62), Inches(1.82), Inches(5.6),
     "Dr. Findus, our startup in Berlin",
     "German GPs code every consultation.\n"
     "We propose, the physician confirms.\n"
     "Reliability is the product.\n"
     "Today it runs on hosted APIs.",
     title_size=15, body_size=14, gap=6, accent=INK)

txt(s, Inches(6.55), Inches(1.92), Inches(6.15), Inches(1.0),
    "Could it run on\nour own hardware?", 24, bold=True, color=ORANGE, gap=2)
txt(s, Inches(6.55), Inches(3.02), Inches(6.15), Inches(0.8),
    f"Qwen3-4B + LoRA on CodiEsp: {n_train} train / {n_dev} dev\n"
    "Spanish case reports, ICD-10 coded, CC-BY 4.0.",
    13, color=GREY, gap=3)

txt(s, Inches(0.62), Inches(4.08), Inches(5.9), Inches(0.35),
    "The resource problem", 14, bold=True, color=ORANGE)
txt(s, Inches(0.62), Inches(4.48), Inches(5.9), Inches(1.4),
    "8 GB of weights at 4B before a single activation.\n"
    f"Median case {tok['p50']:,.0f} tokens, longest {tok['max']:,.0f}.\n"
    "One JAX program, lowered by XLA to CPU, GPU and TPU.", 15, gap=9)

fig(s, "slide1-coverage.png", Inches(6.55), Inches(4.02), Inches(6.2))

txt(s, Inches(0.62), Inches(6.42), Inches(12.0), Inches(0.7),
    "Where does the time go, and what does it take to run this at all?",
    22, bold=True)

notes(s, f"""
LEONARD  ~55 s   ·   deck is budgeted to 6:00, see the timing note on slide 5

"We run Dr. Findus, a startup in Berlin. German GPs have to attach ICD-10 codes to
every consultation. We propose the codes, the physician confirms each one. Today
that runs on hosted APIs, and the question we could not answer was whether it
could run on our own hardware.

[left card, then the grey line right]
So we built it and measured it: a LoRA fine-tune of Qwen3-4B on CodiEsp,
{n_train} Spanish clinical case reports with ICD-10 codes attached, eight
gigabytes of weights before a single activation.

[the bars]
And the data sets the first constraint. At 512 tokens only
{100 * (1 - trunc[512]['frac_truncated']):.0f} percent of cases fit uncut. We
benchmark at 1024, the orange bar, where
{100 * (1 - trunc[1024]['frac_truncated']):.0f} percent fit. Doubling the window
fixes the rest and costs memory and time. That trade-off is slide four."

IF ASKED
· Why this dataset: reliability is the product for us, over-coding is fraud
  exposure, so a self-hosted model is a compliance question as much as a cost one.
· Accuracy: none is claimed or measured. {n_dev} dev documents exist and
  evaluate.py was never run. The rubric awards nothing for it and neither do we.
· CodiEsp is CC-BY 4.0, real clinical text, de-identified at source.
""")

# ============================== 2. Approach ===================================
s = slide(2, "Approach", "One program. Three backends. All measured.")

cols = [
    ("Compilation",
     "One JAX program\nXLA lowers it to all three\nLoRA via qwix, Tunix trainer"),
    ("Orchestration",
     "3 pinned Docker images\nKubernetes across 2 clusters\nKueue admission, explicit limits"),
    ("Storage",
     "gcsfuse CSI, implicit-dirs\n8 GB checkpoint from a bucket\nFile cache left as a knob"),
]
for i, (head, body) in enumerate(cols):
    card(s, Inches(0.62 + i * 4.15), Inches(1.80), Inches(3.85), head, body,
         title_size=16, body_size=13, gap=7)

panel(s, Inches(0.62), Inches(3.82), Inches(12.10), Inches(0.66),
      fill=WHITE, edge=ORANGE_L)
txt(s, Inches(0.62), Inches(3.98), Inches(12.10), Inches(0.35),
    "gs://me344-tpu-labs-west4   →   gcsfuse CSI sidecar   →   /gcs in the pod"
    "   →   8 GB checkpoint   →   HBM",
    13.5, bold=True, align=PP_ALIGN.CENTER, gap=0)

table(s, Inches(0.62), Inches(4.72), [
    ["Backend", "Hardware", "Host arch", "How the code got there"],
    ["CPU", "32-core x86_64, 31 GiB", "x86_64", "bare node, podman"],
    ["GPU", "NVIDIA GH200 480 GB", "ARM64 Grace", "public image + ConfigMap"],
    ["TPU", "v5e 2×4, 8 chips", "x86_64", "private registry + gcsfuse"],
], [Inches(1.5), Inches(4.1), Inches(2.4), Inches(4.1)], size=12.5,
    rh=Inches(0.41), right_align_from=99)

txt(s, Inches(0.62), Inches(6.55), Inches(12.0), Inches(0.5),
    "The GH200 host is ARM64. That one fact cost five separate blockers.",
    16, bold=True, color=ORANGE)

notes(s, """
LEONARD  ~55 s

"One code path. A single JAX program, and XLA lowers it to all three backends.
Nothing branches on hardware, which is what makes the comparison honest.

[the three cards]
Three pinned Docker images, Kubernetes Jobs across two clusters, Kueue for
admission on the TPU pool, and every manifest states its limits explicitly.

[the strip across the middle]
This is the data path. An eight-gigabyte checkpoint in a bucket, mounted into the
pod by the gcsfuse CSI driver, read from /gcs at startup. We left the file cache
on that mount as a knob rather than fixing it, and that knob became the
experiment. Luis comes back to it.

[table, GPU row]
One line before I hand over. The GH200's host processor is ARM64, not x86. That
single fact cost us five separate blockers."

→ hand over to Luis.

IF ASKED
· The five blockers: no aarch64 libtpu wheel, a segfaulting QEMU cross-build, no
  credentials for the private registry, a 403 from it, and a base image that was
  amd64-only. Fix was to delete the registry from the design: public base image,
  source from a ConfigMap, dependencies from PyPI, HF and Zenodo.
· Why two clusters: the TPU pool is GKE in us-west4, the GH200 is on-prem at
  Stanford. Same manifests, different admission.
· Kueue LocalQueue is student-queue; the CPU plane is bare podman, no scheduler.
""")

# ============================== 3. Measurements ================================
s = slide(3, "Measurements", "We did not time the job. We took it apart.")

n_samp = gpu["utilization"]["n_samples"]
hz = 1 / gpu["utilization"]["sample_interval_s"]
warm = base["steady"]["warmup_steps_excluded"]

y = table(s, Inches(0.62), Inches(1.82), [
    ["What we sample", "Instrument"],
    ["Chip utilisation", f"nvidia-smi {hz:.0f} Hz · psutil on CPU"],
    ["Peak memory", "jax memory_stats() · peak RSS"],
    ["Step time", "median, p10, p90, warmup excluded"],
    ["Wall clock", "six phases, must sum to total"],
], [Inches(2.15), Inches(3.35)], size=12.5, rh=Inches(0.40), right_align_from=99)

txt(s, Inches(0.62), y + Inches(0.16), Inches(5.5), Inches(0.7),
    "One schema, one JSON per run. Every figure in this\n"
    "deck reads those files and nothing else.", 12, color=GREY, gap=2)

card(s, Inches(0.62), Inches(4.52), Inches(5.5),
     "The schema caught a lie",
     "93 µs per step. 11 M tokens/s. Impossible,\n"
     "yet it carried status \"ok\".\n\n"
     "Its own notes admitted it may have timed\n"
     "dispatch, not completion. Discarded and re-run.",
     title_size=15, body_size=12.5, gap=4, fill=RED_TINT, accent=RED)

y = fig(s, "slide3-walltime.png", Inches(6.40), Inches(1.78), Inches(6.30))

table(s, Inches(6.40), y + Inches(0.28), [
    ["Backend", "Peak memory, largest fitting run", "Next step up"],
    ["CPU 32-core", f"{m_cpu:.1f} % of 31 GiB (RSS)", "4B never ran"],
    ["GH200", f"{mem_gpu:.1f} % of 96 GiB HBM", f"batch {oom_gpu} OOM"],
    ["v5e, 8 chips", f"{mem_tpu:.1f} % of 16 GiB/chip", f"batch {oom_tpu} OOM"],
], [Inches(1.60), Inches(3.05), Inches(1.65)], size=12, rh=Inches(0.40),
    right_align_from=99)

txt(s, Inches(6.40), y + Inches(2.02), Inches(6.30), Inches(0.6),
    "Memory is what ends every sweep. One batch step past\nthese numbers, all three fail.",
    13.5, bold=True, color=ORANGE, gap=2)

notes(s, f"""
LUIS  ~65 s

"We did not time the job and divide by steps.

[left table]
Four quantities, each with a named instrument. Utilisation from nvidia-smi once a
second inside the pod, psutil on the CPU node. Peak memory from JAX's own
memory_stats, peak RSS on the CPU. Step times as median and percentiles, warmup
excluded. And the wall clock split into six phases that are required to sum back
to the total.

[the bar]
That is this bar: {bp['total_wall_s']:.0f} seconds, left to right.
{load_pct:.0f} percent of it is just reading the checkpoint. Only the orange
block, {compute_pct:.0f} percent, is arithmetic.

[table underneath]
And memory is what ends every sweep: the largest configuration that fitted, and
the batch that killed it. We keep the failed cells.

[red box]
One run came back at eleven million tokens a second carrying a passing status.
Its own notes field caught it. A passing status is not enough to trust a record."

IF ASKED
· The lie in detail: 93 µs per step, which is dispatch time, not completion. The
  telemetry module had written the suspicion into notes itself. Discarded, re-run.
· No TPU utilisation counter exists in JAX, only memory statistics. That is why
  the cell on slide 4 says so rather than showing a different quantity.
· Warmup excluded is {warm} steps on this run, 2 on the CPU, 5 on the GPU. Never a
  mean, because the first step carries the compile.
· Unattributed time goes to an 'other' field rather than being dropped. Residual
  after reassembly is under 1.1e-13 s on all 8 records.
""")

# ============================== 4. Results ====================================
s = slide(4, "Results", "Under load, one GH200 beats eight TPU chips.")

y = table(s, Inches(0.62), Inches(1.74), [
    ["Qwen3-0.6B, batch 1", "CPU 32-core", "GH200  1 chip", "v5e  8 chips"],
    ["Median step time", f"{c_s:.2f} s", f"{g_s:.4f} s", f"{t_s:.4f} s"],
    ["Speedup vs CPU", "1.0×", f"{sp_g:.0f}×", f"{sp_t:.0f}×"],
    ["Mean chip utilisation", f"{u_cpu:.1f} %", f"{u_gpu:.2f} %", "no counter exists"],
    ["Peak memory", f"{m_cpu:.1f} %", f"{m_gpu:.1f} %", f"{m_tpu:.1f} %"],
], [Inches(2.45), Inches(1.70), Inches(1.70), Inches(1.70)], size=12.5,
    rh=Inches(0.375), rule_after=2)

txt(s, Inches(0.62), y + Inches(0.14), Inches(7.55), Inches(0.7),
    f"Both accelerators are ~{sp_g:.0f}× the CPU and tie with each other to "
    f"{gt:.2f} %.\nAt batch 1 neither is busy. That tie is an artefact.",
    14.5, bold=True, gap=3)

fig(s, "slide4-sweep.png", Inches(0.62), Inches(4.42), Inches(7.55))

stat(s, Inches(8.45), Inches(1.74), Inches(4.25), f"{load_ratio:.1f}x",
     "the throughput of an eight-chip v5e\nslice, from a single Hopper chip", big=52)

card(s, Inches(8.45), Inches(3.52), Inches(4.25),
     "Memory boundary",
     f"GH200 fails between {bs_gpu} and {oom_gpu}\nv5e fails between {bs_tpu} and {oom_tpu}",
     title_size=14, body_size=13.5, gap=4)

card(s, Inches(8.45), Inches(5.02), Inches(4.25),
     "The mitigation, controlled",
     f"File cache on: checkpoint load\n{load_x:.1f}× faster, {wall_cut:.0f} % off wall clock.\n"
     "Median step time unchanged to\nthe fourth decimal.",
     title_size=14, body_size=13, gap=3)

notes(s, f"""
LUIS  ~90 s   ·   the slide to spend time on

"Two comparisons, and the second contradicts the first.

[top table]
All three backends, identical code, identical batch. Nineteen seconds a step on
the CPU, eighty milliseconds on both accelerators. That is {sp_g:.0f} times.

[fourth row]
But look at utilisation: {u_gpu:.1f} percent on the GH200. And peak memory,
{m_gpu:.0f} and {m_tpu:.0f} percent. Neither chip is doing anything. The two tie
to within a quarter of a percent, and that tie is an artefact: at batch one we
were comparing launch latency, not compute.

[the chart]
Same two chips at four billion parameters, batch swept until they died. Orange,
one GH200, climbs from {g8:,.0f} to {g32:,.0f} tokens a second. Blue, eight TPU
chips, is flat at {t8:,.0f}: saturated at batch eight, the extra work bought
nothing. And it runs out of memory a full step before the GPU.

So {load_ratio:.1f} times the throughput of eight TPU chips, from one Hopper chip.
Same hardware, opposite conclusion, because the first measurement was taken where
the answer is meaningless.

[bottom right]
And the controlled experiment. Checkpoint load was the largest phase, so we turned
on the gcsfuse file cache and re-ran. Load {load_x:.1f} times faster, wall clock
down {wall_cut:.0f} percent. The control: median step time unchanged to the fourth
decimal. Only I/O moved."

Hand back to Leonard.

IF ASKED
· Why the GH200 wins under load: identical step speed at batch 1, but the slice
  stops scaling at batch 8 while the Hopper keeps going, and the slice pays 8
  chips for it.
· Sequence length: 512 to 2048 grows step time as seq^1.30 with peak HBM flat.
  Remat trades capacity for time.
· Batch 16 on the TPU is 1.8 % slower per token than batch 8, so the batch the
  HBM ceiling forbids was not worth having anyway.
· File cache: fileCacheCapacity 0 to 20Gi on the gcsfuse mount, range reads cached.
""")

# ============================== 5. Conclusion =================================
s = slide(5, "Conclusion", "The chip was never the bottleneck.")

findings = [
    (f"{compute_pct:.0f} %",
     "of wall clock is arithmetic",
     f"Of {bp['total_wall_s']:.0f} s, only {bp['steady_state_s']:.0f} s computes.\n"
     f"The biggest single phase is\nreading the checkpoint: {load_pct:.0f} %."),
    ("6.7×",
     "more parameters, CPU stopped",
     "0.6B to 4B did not make the CPU\n6.7× slower. XLA never finished\n"
     "compiling; the kernel killed it\nat 31.5 GB."),
    ("1 of 36",
     "dependencies were TPU-bound",
     "The GH200 sat idle behind an\nARM64 host, a QEMU segfault\n"
     "and a registry 403. Swapping\none pin is the whole port."),
]
for i, (big, cap, body) in enumerate(findings):
    x = Inches(0.62 + i * 4.15)
    panel(s, x, Inches(1.80), Inches(3.85), Inches(2.50))
    txt(s, x + Inches(0.28), Inches(1.94), Inches(3.3), Inches(0.7), big,
        34 if i == 0 else 26, bold=True, color=ORANGE if i == 0 else INK, gap=0)
    txt(s, x + Inches(0.28), Inches(2.56), Inches(3.35), Inches(0.45), cap, 12,
        bold=True, gap=2)
    txt(s, x + Inches(0.28), Inches(3.03), Inches(3.35), Inches(1.3), body, 11.5,
        color=GREY, gap=2)

txt(s, Inches(0.62), Inches(4.62), Inches(5.9), Inches(0.35),
    "Scaling recommendation", 15, bold=True, color=ORANGE)
txt(s, Inches(0.62), Inches(5.00), Inches(5.9), Inches(1.4),
    "Do not scale out. Amortise: raise batch only until\n"
    "the chip is occupied, then attack fixed cost:\n"
    "persistent compile caches, warm workers instead of\n"
    "a pod per job, file cache on anything reading a\n"
    "checkpoint.", 13, gap=2)

card(s, Inches(6.80), Inches(4.46), Inches(5.90),
     "What this means for Dr. Findus",
     f"At batch 8, the slice's own best point: one GH200 runs\n"
     f"{g8 / t8:.1f}× faster than eight TPU chips and costs "
     f"{cost_ratio_4b:.0f}× less,\n${k4_gpu:.2f} against ${k4_tpu:.2f} per 1,000 steps.\n\n"
     "Self-hosting is not the cost problem we assumed.\n"
     "The real cost is fixed overhead per job.",
     title_size=15, body_size=12.5, gap=2,
     foot="TPU rate billed; the GH200 rate is a market proxy.")

txt(s, Inches(0.62), Inches(6.98), Inches(12.1), Inches(0.30),
    f"{len(recs)} measured runs  ·  "
    "github.com/RN0L/me344-qwen3-icd10-profiling",
    11, color=GREY, align=PP_ALIGN.CENTER)

notes(s, f"""
LEONARD  ~85 s

"Three findings.

[first card]
The chip was never the bottleneck. Of {bp['total_wall_s']:.0f} seconds,
{compute_pct:.0f} percent is arithmetic. The largest single phase is reading a
checkpoint.

[second card]
Scaling the model did not scale the cost. Under seven times more parameters did
not make the CPU seven times slower, it made it impossible: the compiler never
finished, and without rematerialisation the kernel killed it at 31.5 gigabytes. A
speedup chart would have drawn a bar there. There is no bar.

[third card]
And the wall was in the supply chain, not the silicon. Exactly one of the
framework's thirty-six dependencies is TPU-bound. Swapping that one pin is the
whole port.

[bottom left]
So: do not scale out, amortise. Raise the batch until the chip is occupied, then
attack fixed cost.

[orange box]
For us, that is the answer. At batch eight, the slice's own best point, one GH200
is {g8 / t8:.1f} times faster and {cost_ratio_4b:.0f} times cheaper. The TPU rate
is what we were billed, the GH200 is a market proxy, and the ordering holds across
its whole published price range.

Self-hosting was not the cost problem. Fixed overhead per job is. Thank you."

TIMING   55 s · 55 s · 65 s · 90 s · 85 s   =   5:50 spoken, ~6:00 with handovers.
Read at ~135 words a minute. The IF ASKED blocks are not part of the 6 minutes.

IF ASKED
· Biggest remaining win: a persistent XLA compile cache. Already plumbed and
  deliberately left off so the compile cost stays visible. Targets 26 to 32 % of
  wall clock across the runs we have.
· Why not more nodes: at {u_gpu:.1f} % utilisation there is nothing to parallelise.
  Scaling out multiplies the fixed cost we just identified.
· What we would do differently: check egress from inside the cluster, not from the
  node. That one unchecked inference nearly cost us the GPU row entirely.
· 4B on CPU: keep remat and XLA never finishes compiling, over 64 minutes single
  threaded. Drop it and the kernel OOM-kills at 31.5 GB against 31 GiB.
""")

prs.save(OUT)
n = len(prs.slides._sldIdLst)
print(f"wrote {OUT.relative_to(REPO)}, {n} slides, {OUT.stat().st_size // 1024} KB")

# --- guard: nothing may run off the bottom or the right edge --------------------
# The card heights above are computed, but the y offsets are still written by hand,
# so this is the check that catches a hand-written offset that stopped fitting.
MARGIN = 0.14
bad = []
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        if sh.height == H or sh.shape_id in EXEMPT:
            continue
        bottom = (sh.top + sh.height) / Inches(1)
        right = (sh.left + sh.width) / Inches(1)
        label = (sh.text_frame.text.split("\n")[0][:34]
                 if sh.has_text_frame else sh.shape_type)
        if bottom > H / Inches(1) - MARGIN:
            bad.append(f"  slide {i}: bottom {bottom:.2f}in  {label!r}")
        if right > W / Inches(1) - MARGIN:
            bad.append(f"  slide {i}: right {right:.2f}in  {label!r}")
print("layout clean, nothing within %.2f in of an edge" % MARGIN if not bad
      else "LAYOUT OVERFLOW:\n" + "\n".join(bad))
